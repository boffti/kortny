#!/usr/bin/env python3
"""Build a .pptx deck from a JSON spec — runs in the Kortny sandbox.

The spec is a list of slides; each slide names a layout and supplies the
fields that layout needs. The script enforces the deck conventions Kortny
commits to:

  * one idea per slide — titles are sentences, not labels
  * a consistent theme (single accent color, two font sizes) applied to
    every slide so the deck reads as one document, not a pile of templates
  * generous margins; body text never runs edge-to-edge
  * an optional source line anchored to the slide footer for any data slide

Network is never required. Input is a JSON file (``--spec``) or stdin;
output is a .pptx path under the workbench dir.

Supported layouts (see references/spec-format.md for the full reference):

  * ``title``   — title + subtitle (deck cover)
  * ``section`` — section divider (big centered title)
  * ``bullets`` — title + bullet list (max 6 bullets enforced)
  * ``two_col`` — title + left/right bullet columns
  * ``quote``   — large centered pull-quote + attribution

Spec shape:

    {
      "theme": {"accent": "2563EB", "title_pt": 36, "body_pt": 18},
      "slides": [
        {"layout": "title", "title": "Q1 Review", "subtitle": "Kortny"},
        {"layout": "bullets", "title": "What moved",
         "bullets": ["Signups +40%", "Churn flat"], "source": "Mixpanel, Apr 1"}
      ]
    }
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

EMU_SLIDE_W = Inches(13.333)  # 16:9 widescreen
EMU_SLIDE_H = Inches(7.5)
MARGIN = Inches(0.9)
MAX_BULLETS = 6

DEFAULT_ACCENT = "2563EB"
DEFAULT_TITLE_PT = 36
DEFAULT_BODY_PT = 18
INK = RGBColor(0x1F, 0x29, 0x37)  # near-black body text
MUTED = RGBColor(0x6B, 0x72, 0x80)  # footer / source line


def _accent(theme: dict[str, Any]) -> RGBColor:
    hex_value = str(theme.get("accent", DEFAULT_ACCENT)).lstrip("#")
    return RGBColor.from_string(hex_value.upper())


def _add_textbox(slide: Any, left: Any, top: Any, width: Any, height: Any) -> Any:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _set_title(
    slide: Any, text: str, theme: dict[str, Any], *, centered: bool = False
) -> None:
    tf = _add_textbox(slide, MARGIN, Inches(0.6), EMU_SLIDE_W - 2 * MARGIN, Inches(1.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(int(theme.get("title_pt", DEFAULT_TITLE_PT)))
    run.font.bold = True
    run.font.color.rgb = _accent(theme)


def _set_source(slide: Any, source: str | None, theme: dict[str, Any]) -> None:
    if not source:
        return
    tf = _add_textbox(
        slide, MARGIN, EMU_SLIDE_H - Inches(0.6), EMU_SLIDE_W - 2 * MARGIN, Inches(0.4)
    )
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Source: {source}"
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = MUTED


def _bullets_into(tf: Any, bullets: list[str], theme: dict[str, Any]) -> None:
    body_pt = int(theme.get("body_pt", DEFAULT_BODY_PT))
    for i, text in enumerate(bullets[:MAX_BULLETS]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {text}"
        run.font.size = Pt(body_pt)
        run.font.color.rgb = INK


def _render_slide(
    prs: Presentation, slide_spec: dict[str, Any], theme: dict[str, Any]
) -> None:
    layout = slide_spec.get("layout", "bullets")
    blank = prs.slide_layouts[6]  # blank layout — we place everything by hand
    slide = prs.slides.add_slide(blank)

    if layout == "title":
        _set_title(slide, slide_spec.get("title", ""), theme, centered=True)
        subtitle = slide_spec.get("subtitle")
        if subtitle:
            tf = _add_textbox(
                slide, MARGIN, Inches(2.4), EMU_SLIDE_W - 2 * MARGIN, Inches(1.0)
            )
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = subtitle
            run.font.size = Pt(int(theme.get("body_pt", DEFAULT_BODY_PT)) + 4)
            run.font.color.rgb = MUTED
    elif layout == "section":
        tf = _add_textbox(
            slide, MARGIN, Inches(3.0), EMU_SLIDE_W - 2 * MARGIN, Inches(1.5)
        )
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slide_spec.get("title", "")
        run.font.size = Pt(int(theme.get("title_pt", DEFAULT_TITLE_PT)) + 8)
        run.font.bold = True
        run.font.color.rgb = _accent(theme)
    elif layout == "bullets":
        _set_title(slide, slide_spec.get("title", ""), theme)
        tf = _add_textbox(
            slide, MARGIN, Inches(2.0), EMU_SLIDE_W - 2 * MARGIN, Inches(4.5)
        )
        _bullets_into(tf, list(slide_spec.get("bullets", [])), theme)
        _set_source(slide, slide_spec.get("source"), theme)
    elif layout == "two_col":
        _set_title(slide, slide_spec.get("title", ""), theme)
        half = (EMU_SLIDE_W - 2 * MARGIN - Inches(0.5)) / 2
        left_tf = _add_textbox(slide, MARGIN, Inches(2.0), half, Inches(4.5))
        _bullets_into(left_tf, list(slide_spec.get("left", [])), theme)
        right_tf = _add_textbox(
            slide, MARGIN + half + Inches(0.5), Inches(2.0), half, Inches(4.5)
        )
        _bullets_into(right_tf, list(slide_spec.get("right", [])), theme)
        _set_source(slide, slide_spec.get("source"), theme)
    elif layout == "quote":
        tf = _add_textbox(
            slide, MARGIN, Inches(2.4), EMU_SLIDE_W - 2 * MARGIN, Inches(2.5)
        )
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"“{slide_spec.get('quote', '')}”"
        run.font.size = Pt(int(theme.get("title_pt", DEFAULT_TITLE_PT)) - 2)
        run.font.italic = True
        run.font.color.rgb = INK
        attribution = slide_spec.get("attribution")
        if attribution:
            ap = tf.add_paragraph()
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            ar.text = f"— {attribution}"
            ar.font.size = Pt(int(theme.get("body_pt", DEFAULT_BODY_PT)))
            ar.font.color.rgb = MUTED
    else:
        raise ValueError(f"unknown layout {layout!r}")


def build(spec: dict[str, Any], out_path: Path) -> Path:
    slides = spec.get("slides")
    if not slides:
        raise ValueError("spec must contain at least one slide under 'slides'")
    theme = dict(spec.get("theme", {}))
    prs = Presentation()
    prs.slide_width = EMU_SLIDE_W
    prs.slide_height = EMU_SLIDE_H
    for slide_spec in slides:
        _render_slide(prs, slide_spec, theme)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Build a .pptx deck from a JSON spec.")
    parser.add_argument("--spec", help="path to a JSON spec file; omit to read stdin")
    parser.add_argument("--out", required=True, help="output .pptx path")
    args = parser.parse_args(argv)

    raw = Path(args.spec).read_text() if args.spec else sys.stdin.read()
    spec = json.loads(raw)
    out = build(spec, Path(args.out))
    print(f"wrote deck: {out} ({out.stat().st_size} bytes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
